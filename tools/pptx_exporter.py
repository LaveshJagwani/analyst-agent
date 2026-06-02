from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
from design_systems import DESIGN_SYSTEMS

def hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert hex string (e.g. #0F172A) to RGBColor object."""
    hex_str = hex_str.lstrip('#')
    r, g, b = tuple(int(hex_str[i:i+2], 16) for i in (0, 2, 4))
    return RGBColor(r, g, b)

def apply_design(slide, design_spec: dict):
    """Apply background and fonts based on design system."""
    system_name = design_spec.get("design_system", "modern_blue")
    system = DESIGN_SYSTEMS.get(system_name, DESIGN_SYSTEMS["modern_blue"])
    palette = system["palette"]
    
    bg_color = hex_to_rgb(palette[0])
    
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    return system

# ── Renderers ─────────────────────────────────────────────────────────────────

def render_title_slide(prs, slide_data: dict, design_spec: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    system = apply_design(slide, design_spec)
    palette = system["palette"]
    font_name = system["font"]
    
    primary_color = hex_to_rgb(palette[2])
    accent_color = hex_to_rgb(palette[1])
    
    # Title
    t = slide.shapes.title
    t.text = slide_data.get("title", "Analysis Report")
    t.text_frame.paragraphs[0].font.color.rgb = primary_color
    t.text_frame.paragraphs[0].font.name = font_name
    t.text_frame.paragraphs[0].font.size = Pt(44)
    t.text_frame.paragraphs[0].font.bold = True
    
    # Subtitle
    if len(slide.placeholders) > 1:
        s = slide.placeholders[1]
        s.text = slide_data.get("subtitle", "")
        s.text_frame.paragraphs[0].font.color.rgb = accent_color
        s.text_frame.paragraphs[0].font.name = font_name
        s.text_frame.paragraphs[0].font.size = Pt(24)

def render_headline_metric(prs, slide_data: dict, design_spec: dict):
    """One big metric and a headline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    system = apply_design(slide, design_spec)
    palette = system["palette"]
    font_name = system["font"]
    accent_color = hex_to_rgb(palette[1])
    text_color = hex_to_rgb(palette[2])
    
    # Headline
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data.get("title", "")
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = text_color
    p.font.name = font_name
    
    # Metric (Large)
    content = slide_data.get("content", [])
    metric_text = content[0] if content else "N/A"
    
    mBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    mtf = mBox.text_frame
    mp = mtf.paragraphs[0]
    mp.text = metric_text
    mp.font.bold = True
    mp.font.size = Pt(72)
    mp.font.color.rgb = accent_color
    mp.font.name = font_name
    mp.alignment = PP_ALIGN.CENTER

def render_chart_left_text_right(prs, slide_data: dict, design_spec: dict, charts_dir: Path):
    """Chart on left, bullets on right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    system = apply_design(slide, design_spec)
    palette = system["palette"]
    font_name = system["font"]
    text_color = hex_to_rgb(palette[2])
    
    # Title
    tBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(1))
    p = tBox.text_frame.paragraphs[0]
    p.text = slide_data.get("title", "")
    p.font.bold = True; p.font.size = Pt(28); p.font.color.rgb = text_color; p.font.name = font_name
    
    # Chart (Left)
    c_ref = slide_data.get("chart_reference")
    chart_placed = False
    if c_ref and charts_dir:
        p = charts_dir / Path(c_ref).name
        if p.exists():
            slide.shapes.add_picture(str(p), Inches(0.5), Inches(1.6), Inches(6.5), Inches(5.0))
            chart_placed = True
            
    # Text (Right)
    bx = Inches(7.5) if chart_placed else Inches(1.0)
    bw = Inches(5.0) if chart_placed else Inches(11.0)
    
    bBox = slide.shapes.add_textbox(bx, Inches(1.6), bw, Inches(5.0))
    tf = bBox.text_frame
    tf.word_wrap = True
    for bullet in slide_data.get("content", []):
        p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18); p.font.color.rgb = text_color; p.font.name = font_name
        p.space_after = Pt(12)

def render_comparison_blocks(prs, slide_data: dict, design_spec: dict):
    """Side by side blocks."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    system = apply_design(slide, design_spec)
    palette = system["palette"]
    font_name = system["font"]
    text_color = hex_to_rgb(palette[2])
    accent_color = hex_to_rgb(palette[1])
    
    # Title
    tBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(1))
    p = tBox.text_frame.paragraphs[0]
    p.text = slide_data.get("title", "")
    p.font.bold = True; p.font.size = Pt(28); p.font.color.rgb = accent_color; p.font.name = font_name
    
    # Content Blocks
    content = slide_data.get("content", [])
    for i, item in enumerate(content[:3]): # Max 3 blocks
        x = Inches(0.5 + i*4.2)
        box = slide.shapes.add_textbox(x, Inches(2.0), Inches(4.0), Inches(4.5))
        box.text_frame.word_wrap = True
        p = box.text_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(16); p.font.color.rgb = text_color; p.font.name = font_name

def export_pptx(result: dict, output_path: Path, charts_dir: Path = None):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    payload = result.get("presentation", {})
    slides = payload.get("slides", [])
    design_spec = payload.get("design", {}) or {
        "design_system": "modern_blue",
        "layout_map": {},
        "visual_rules": {}
    }
    
    layout_map = design_spec.get("layout_map", {})
    
    for i, slide_data in enumerate(slides):
        stype = slide_data.get("type", "generic")
        
        # Decide layout
        layout = "generic"
        if stype == "title":
            layout = "title"
        elif stype == "summary":
            layout = layout_map.get("executive_summary", "headline_metric")
        elif stype == "chart":
            layout = layout_map.get("main_insights", "chart_left_text_right")
        elif stype == "recommendations":
            layout = layout_map.get("recommendations", "comparison_blocks")
            
        try:
            if layout == "title":
                render_title_slide(prs, slide_data, design_spec)
            elif layout == "headline_metric":
                render_headline_metric(prs, slide_data, design_spec)
            elif layout == "chart_left_text_right":
                render_chart_left_text_right(prs, slide_data, design_spec, charts_dir)
            elif layout == "comparison_blocks":
                render_comparison_blocks(prs, slide_data, design_spec)
            else:
                # Default generic
                render_chart_left_text_right(prs, slide_data, design_spec, charts_dir)
        except Exception as e:
            print(f"Error rendering slide {i}: {e}")
            
    prs.save(str(output_path))
    return str(output_path)
